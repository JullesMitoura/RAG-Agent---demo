import tempfile
import os
from langchain_text_splitters import TokenTextSplitter
from qdrant_client.models import Distance, VectorParams
from langchain_qdrant import Qdrant
from app.services.qdrant_service import QdrantService
from app.services.azure_services import AzureServices
from pptx import Presentation
import docx
import PyPDF2

class RagService:
    def __init__(self):
        self.qdrant_services = QdrantService()
        self.azure_services = AzureServices()

    def docx_read(self, file):
        """Read text from a .docx file."""
        doc = docx.Document(file)
        text = [part.text for part in doc.paragraphs if part.text.strip()]
        return '\n'.join(text)

    def ppt_read(self, file):
        """Read text from a .pptx file."""
        doc = Presentation(file)
        text = []
        for slide in doc.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
        return '\n'.join(text)

    def pdf_read(self, file):
        """Read text from a .pdf file."""
        text = []
        with open(file, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                if page.extract_text():
                    text.append(page.extract_text())
        return '\n'.join(text)

    def read_file(self, file, file_type):
        """Read content from a file based on its type."""
        print('-----')
        print(file_type)
        if file_type == '.docx':
            return self.docx_read(file)
        elif file_type == '.pptx':
            return self.ppt_read(file)
        elif file_type == '.pdf':
            return self.pdf_read(file)
        elif file_type == '.txt':
            with open(file, 'r', encoding='utf-8') as f:
                return f.read()
        else:
            raise ValueError(f"Unsupported file type: {file_type}. Supported types are: docx, pptx, pdf, txt.")

    def doc_chunking(self, temp_file_path, file_type):
        """Chunk the document into smaller parts using a temporary file."""
        content = self.read_file(temp_file_path, file_type)
        text_splitter = TokenTextSplitter(chunk_size=500, chunk_overlap=50)
        texts = text_splitter.split_text(content)
        return texts
    
    def qdrant_processing(self, file_content: bytes, file_name: str, collection_name: str, file_type: str):
        """Process the uploaded file content and add its chunks to Qdrant."""
        client = self.qdrant_services.qdrant_client
        embedding_model = self.azure_services.define_embedding()
        if client.collection_exists(collection_name):
            print(f"Collection '{collection_name}' already exists.")
        else:
            print(f"Creating collection '{collection_name}'.")
            client.create_collection(
                collection_name,
                vectors_config=VectorParams(size=1536, distance=Distance.DOT)
            )

        with tempfile.NamedTemporaryFile(delete=True, mode='wb') as temp_file:
            temp_file.write(file_content)
            temp_file.flush()

            texts = self.doc_chunking(temp_file.name, file_type)
            print(f"Adding {len(texts)} chunks to Qdrant.")

            if texts:
                qdrant = Qdrant(client, collection_name, embedding_model)
                qdrant.add_texts(texts, metadatas=[{'name': file_name}] * len(texts))
            else:
                print("No texts to add to Qdrant.")